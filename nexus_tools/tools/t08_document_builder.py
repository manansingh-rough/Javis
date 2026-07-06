"""
NEXUS AI v4.0 — Tool 08: Document generation.
Hardware: Intel i3 7th Gen · 12GB RAM · Ollama + Groq API

Generates documents in DOCX, PPTX, Markdown, and plain text formats.
Supports headers, paragraphs, lists, tables, images, and basic formatting.
"""

import json
import logging
import time
from pathlib import Path
from typing import Optional, List, Dict, Any

logger = logging.getLogger("nexus.tool.document_builder")


def document_builder(
    format: str,
    title: str,
    content: List[Dict[str, Any]],
    output_path: str,
    author: str = "NEXUS AI",
) -> str:
    """
    Generate formatted documents in various formats.
    
    Use this tool when: The user asks to create a document, report, presentation,
    markdown file, or any formatted document output.
    
    Args:
        format: Output format: "markdown", "docx", "pptx", "text".
        title: Document title (used as filename base and document title).
        content: List of content blocks. Each block is a dict with:
                 - "type": "heading", "paragraph", "code", "list", "table", "image"
                 - "level": heading level (1-6, for "heading" type)
                 - "text": text content
                 - "items": list of strings (for "list" type)
                 - "rows": list of lists (for "table" type)
                 - "headers": list of strings (for "table" type)
                 - "code": code string (for "code" type)
                 - "language": language name (for "code" type)
                 - "path": image path (for "image" type)
        output_path: Path where the document will be saved.
        author: Document author name.
    
    Returns:
        JSON string with keys:
          - success (bool): Whether the document was generated.
          - result (str): Path to the generated file.
          - error (str or null): Error message if failed.
    
    Examples:
        >>> document_builder("markdown", "My Report", [
        ...     {"type": "heading", "text": "Introduction", "level": 1},
        ...     {"type": "paragraph", "text": "This is a report."},
        ... ], "output/report.md")
    """
    start = time.perf_counter()
    
    try:
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        
        if format == "markdown":
            return _build_markdown(title, content, out)
        elif format == "text":
            return _build_text(title, content, out)
        elif format == "docx":
            return _build_docx(title, content, out, author)
        elif format == "pptx":
            return _build_pptx(title, content, out, author)
        else:
            return json.dumps({
                "success": False, "result": None,
                "error": f"Unknown format: '{format}'. Valid: markdown, docx, pptx, text"
            })
    
    except Exception as e:
        logger.error(f"document_builder error: {e}", exc_info=True)
        return json.dumps({
            "success": False, "result": None,
            "error": f"{type(e).__name__}: {e}"
        })


def _build_markdown(title: str, content: List[Dict], out: Path) -> str:
    """Build a Markdown document."""
    lines = [f"# {title}", ""]
    
    for block in content:
        block_type = block.get("type", "paragraph")
        
        if block_type == "heading":
            level = block.get("level", 2)
            lines.append(f"{'#' * level} {block.get('text', '')}")
        elif block_type == "paragraph":
            lines.append(block.get("text", ""))
        elif block_type == "code":
            lang = block.get("language", "")
            lines.append(f"```{lang}")
            lines.append(block.get("code", ""))
            lines.append("```")
        elif block_type == "list":
            for item in block.get("items", []):
                lines.append(f"- {item}")
        elif block_type == "table":
            headers = block.get("headers", [])
            rows = block.get("rows", [])
            if headers:
                lines.append("| " + " | ".join(headers) + " |")
                lines.append("|" + "|".join("---" for _ in headers) + "|")
            for row in rows:
                lines.append("| " + " | ".join(str(c) for c in row) + " |")
        elif block_type == "image":
            lines.append(f"![{block.get('text', 'image')}]({block.get('path', '')})")
        
        lines.append("")
    
    out.write_text("\n".join(lines), encoding="utf-8")
    return json.dumps({"success": True, "result": str(out), "error": None})


def _build_text(title: str, content: List[Dict], out: Path) -> str:
    """Build a plain text document."""
    lines = [title, "=" * len(title), ""]
    
    for block in content:
        block_type = block.get("type", "paragraph")
        
        if block_type == "heading":
            level = block.get("level", 2)
            prefix = "#" * level
            lines.append(f"{prefix} {block.get('text', '')}")
        elif block_type == "paragraph":
            lines.append(block.get("text", ""))
        elif block_type == "code":
            lines.append(block.get("code", ""))
        elif block_type == "list":
            for item in block.get("items", []):
                lines.append(f"  - {item}")
        elif block_type == "table":
            rows = block.get("rows", [])
            for row in rows:
                lines.append("  " + " | ".join(str(c) for c in row))
        
        lines.append("")
    
    out.write_text("\n".join(lines), encoding="utf-8")
    return json.dumps({"success": True, "result": str(out), "error": None})


def _build_docx(title: str, content: List[Dict], out: Path, author: str) -> str:
    """Build a DOCX document using python-docx."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document()
        doc.add_heading(title, 0)
        core_properties = doc.core_properties
        core_properties.author = author
        core_properties.title = title
        
        for block in content:
            block_type = block.get("type", "paragraph")
            
            if block_type == "heading":
                level = min(block.get("level", 2), 9)
                doc.add_heading(block.get("text", ""), level)
            
            elif block_type == "paragraph":
                doc.add_paragraph(block.get("text", ""))
            
            elif block_type == "code":
                p = doc.add_paragraph()
                run = p.add_run(block.get("code", ""))
                run.font.name = "Courier New"
                run.font.size = Pt(9)
            
            elif block_type == "list":
                for item in block.get("items", []):
                    doc.add_paragraph(item, style="List Bullet")
            
            elif block_type == "table":
                headers = block.get("headers", [])
                rows = block.get("rows", [])
                table = doc.add_table(rows=1 + len(rows), cols=max(len(headers), len(rows[0]) if rows else 1))
                table.style = "Light Grid Accent 1"
                if headers:
                    for i, h in enumerate(headers):
                        table.rows[0].cells[i].text = str(h)
                for ri, row in enumerate(rows):
                    for ci, cell in enumerate(row):
                        table.rows[ri + 1].cells[ci].text = str(cell)
            
            elif block_type == "image":
                img_path = block.get("path", "")
                if Path(img_path).exists():
                    doc.add_picture(img_path, width=Inches(5))
        
        doc.save(str(out))
        return json.dumps({"success": True, "result": str(out), "error": None})
    
    except ImportError:
        return json.dumps({
            "success": False, "result": None,
            "error": "python-docx not installed. Install with: pip install python-docx"
        })


def _build_pptx(title: str, content: List[Dict], out: Path, author: str) -> str:
    """Build a PPTX presentation using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation()
        
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"Generated by {author}"
        
        for block in content:
            block_type = block.get("type", "paragraph")
            
            if block_type == "heading":
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = block.get("text", "")
            
            elif block_type == "paragraph":
                slide_layout = prs.slide_layouts[2]  # Content only
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = "Content"
                text_frame = slide.placeholders[1].text_frame
                text_frame.text = block.get("text", "")
            
            elif block_type == "list":
                slide_layout = prs.slide_layouts[2]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = "List"
                text_frame = slide.placeholders[1].text_frame
                for item in block.get("items", []):
                    p = text_frame.add_paragraph()
                    p.text = f"• {item}"
            
            elif block_type == "table":
                rows = block.get("rows", [])
                if rows:
                    slide_layout = prs.slide_layouts[6]  # Blank
                    slide = prs.slides.add_slide(slide_layout)
                    rows_count = len(rows) + 1
                    cols_count = max(len(r) for r in rows)
                    table = slide.shapes.add_table(rows_count, cols_count, Inches(1), Inches(1), Inches(8), Inches(4)).table
                    headers = block.get("headers", [])
                    if headers:
                        for i, h in enumerate(headers):
                            table.cell(0, i).text = str(h)
                    for ri, row in enumerate(rows):
                        for ci, cell in enumerate(row):
                            table.cell(ri + 1, ci).text = str(cell)
        
        prs.save(str(out))
        return json.dumps({"success": True, "result": str(out), "error": None})
    
    except ImportError:
        return json.dumps({
            "success": False, "result": None,
            "error": "python-pptx not installed. Install with: pip install python-pptx"
        })